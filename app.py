import streamlit as st
import pandas as pd
import numpy as np
import re
import holidays
from io import BytesIO
from datetime import datetime, date, timedelta, time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import scipy.stats as stats

# --- 1. CONFIGURACIÓN Y ESTILOS ---
st.set_page_config(page_title="Gestión de Energía - Dashboard SGE", layout="wide", page_icon="🚆")
chile_holidays = holidays.Chile()

ORDEN_TIPO_DIA = ["L", "S", "D/F"]

st.markdown("""
    <style>
    .stMetric { background-color: #ffffff; padding: 20px; border-radius: 10px; border-left: 5px solid #005195; box-shadow: 0 2px 4px rgba(0,0,0,0.05); }
    </style>
    """, unsafe_allow_html=True)

# --- 2. FUNCIONES DE PROCESAMIENTO Y EXPORTACIÓN ---
def to_pptx(title_text, df=None, metrics_dict=None):
    prs = Presentation()
    slide_layout = prs.slide_layouts[5] 
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = f"EFE Valparaíso: {title_text}"
    y_cursor = Inches(1.5)
    if metrics_dict:
        txBox = slide.shapes.add_textbox(Inches(0.5), y_cursor, Inches(9), Inches(1))
        tf = txBox.text_frame
        for k, v in metrics_dict.items():
            p = tf.add_paragraph()
            p.text = f"• {k}: {v}"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 81, 149)
        y_cursor += Inches(1.2)
    if df is not None and not df.empty:
        df_display = df.head(12).reset_index(drop=True)
        rows, cols = df_display.shape
        table = slide.shapes.add_table(rows + 1, cols, Inches(0.5), y_cursor, Inches(9), Inches(3)).table
        for c, col_name in enumerate(df_display.columns):
            cell = table.cell(0, c)
            cell.text = str(col_name)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 81, 149) 
            p = cell.text_frame.paragraphs[0]
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.size = Pt(10)
            p.font.bold = True
        for r in range(rows):
            for c in range(cols):
                val = df_display.iloc[r, c]
                formatted_val = str(val) if not isinstance(val, float) else f"{val:,.1f}"
                table.cell(r + 1, c).text = formatted_val
                table.cell(r + 1, c).text_frame.paragraphs[0].font.size = Pt(9)
    binary_output = BytesIO()
    prs.save(binary_output)
    return binary_output.getvalue()

def exportar_resumen_excel(metrics_dict, df_resumen_jornada, df_energia, df_datos_semanales=None):
    output = BytesIO()
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        df_metrics = pd.DataFrame([metrics_dict]).T.reset_index()
        df_metrics.columns = ['Métrica', 'Valor']
        df_metrics.to_excel(writer, sheet_name='Métricas', index=False)
        if df_resumen_jornada is not None and not df_resumen_jornada.empty:
            df_resumen_jornada.to_excel(writer, sheet_name='Resumen_Jornada', index=False)
        if df_energia is not None and not df_energia.empty:
            df_energia.to_excel(writer, sheet_name='Energía_Prioridad', index=False)
        if df_datos_semanales is not None and not df_datos_semanales.empty:
            df_datos_semanales.to_excel(writer, sheet_name='Datos_Semanales', index=False)
    return output.getvalue()

def parse_latam_number(val):
    if pd.isna(val): return 0.0
    if isinstance(val, (int, float)): return float(val)
    s = str(val).strip().replace(' ', '').replace('$', '')
    s = re.sub(r'[^\d.,-]', '', s)
    if not s: return 0.0
    if ',' in s and '.' in s:
        if s.rfind(',') > s.rfind('.'): s = s.replace('.', '').replace(',', '.')
        else: s = s.replace(',', '')
    elif ',' in s: s = s.replace(',', '.')
    try: return float(s)
    except: return 0.0

def get_tipo_dia(fch):
    if fch in chile_holidays or fch.weekday() == 6: return "D/F"
    if fch.weekday() == 5: return "S"
    return "L"

def to_excel_consolidado(df_ops, df_tr, df_tr_acum, df_seat, df_p_d, df_p_15, df_fact_h, df_fact_d):
    output = BytesIO()
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        dfs = {'Operaciones': df_ops, 'Kms_Diarios_Tren': df_tr, 'Odometros_Acum_Tren': df_tr_acum,
               'SEAT': df_seat, 'PRMTE_D': df_p_d, 'PRMTE_15': df_p_15, 'Fact_H': df_fact_h, 'Fact_D': df_fact_d}
        for name, df in dfs.items():
            if not df.empty: df.to_excel(writer, index=False, sheet_name=name)
    return output.getvalue()

# --- 3. FUNCIONES PARA PROCESAR THDR ---
def convertir_a_minutos(val):
    if pd.isna(val) or str(val).strip() == "":
        return None
    try:
        if isinstance(val, (datetime, time)):
            return val.hour * 60 + val.minute + (val.second / 60.0)
        if isinstance(val, str):
            val = val.strip()
            m_ss = re.search(r'(\d{1,2}):(\d{2}):(\d{2})', val)
            if m_ss:
                return int(m_ss.group(1)) * 60 + int(m_ss.group(2)) + (int(m_ss.group(3)) / 60.0)
            m_mm = re.search(r'(\d{1,2}):(\d{2})', val)
            if m_mm:
                return int(m_mm.group(1)) * 60 + int(m_mm.group(2))
        return None
    except:
        return None

def format_hms(minutos_float, con_signo=False):
    if pd.isna(minutos_float) or minutos_float == 0:
        return "00:00:00"
    signo = ("+" if minutos_float > 0 else "-" if minutos_float < 0 else "") if con_signo else ""
    total_segundos = int(round(abs(minutos_float) * 60))
    h, r = divmod(total_segundos, 3600)
    m, s = divmod(r, 60)
    return f"{signo}{h:02d}:{m:02d}:{s:02d}"

def clasificar_flota_func(motriz):
    try:
        num = int(float(motriz))
        if 1 <= num <= 27:
            return "XT-100"
        if 28 <= num <= 35:
            return "XT-M"
        if 101 <= num <= 110:
            return "SFE (Chino)"
        return "OTRO"
    except:
        return "S/I"

DISTANCIAS = {
    "PU-LI": 43.13, "LI-PU": 43.13, "PU-SA": 29.11, "SA-PU": 29.11,
    "EB-PU": 25.40, "PU-EB": 25.40, "VM-LI": 34.03, "LI-VM": 34.03,
    "VM-PU": 9.10,  "PU-VM": 9.10
}

@st.cache_data
def leer_fecha_archivo(file):
    try:
        df = pd.read_excel(file, nrows=1, header=None)
        val = str(df.iloc[0, 0]).split('.')[0].strip().zfill(6)
        return (int(val[0:2]), int(val[2:4]), 2000 + int(val[4:6]))
    except:
        return None

def procesar_thdr_avanzado(file):
    try:
        # Leer todo el archivo sin encabezado
        try:
            df_raw = pd.read_excel(file, header=None, engine=None)
        except Exception:
            df_raw = pd.read_excel(file, header=None, engine='xlrd')
        
        # Las primeras dos filas son encabezados: fila0 = nombres de estaciones, fila1 = "Hora Llegada"/"Hora Salida"
        header0 = df_raw.iloc[0].fillna('').astype(str)
        header1 = df_raw.iloc[1].fillna('').astype(str)
        column_names = []
        for i in range(len(header0)):
            base = header0[i].strip()
            sub = header1[i].strip()
            if sub in ['Hora Llegada', 'Hora Salida']:
                column_names.append(f"{base}_{sub}")
            else:
                column_names.append(base)
        df = df_raw.iloc[2:].copy()
        df.columns = column_names
        
        # --- Búsqueda flexible de columnas base ---
        def buscar_columna(nombres_posibles):
            for col in df.columns:
                for posible in nombres_posibles:
                    if posible.lower() in col.lower():
                        return col
            return None
        
        # Mapeo de columnas necesarias
        col_recorrido = buscar_columna(['Recorrido', 'Trayecto', 'Ruta'])
        col_servicio = buscar_columna(['Servicio', 'Serv', 'N° Servicio'])
        col_hora_prog = buscar_columna(['Hora_Prog', 'Hora Programada', 'Hora Prog', 'Prog'])
        col_motriz1 = buscar_columna(['Motriz 1', 'Motriz1', 'M1', 'Motor 1'])
        col_motriz2 = buscar_columna(['Motriz 2', 'Motriz2', 'M2', 'Motor 2'])
        col_unidad = buscar_columna(['Unidad', 'Tren', 'Formación'])
        
        # Asignar o crear columnas con valores por defecto
        if col_recorrido:
            df['Recorrido'] = df[col_recorrido]
        else:
            df['Recorrido'] = ''
        if col_servicio:
            df['Servicio'] = df[col_servicio]
        else:
            df['Servicio'] = 0
        if col_hora_prog:
            df['Hora_Prog'] = df[col_hora_prog]
        else:
            df['Hora_Prog'] = '00:00:00'
        
        # Motriz 1: si no se encuentra, crear Serie de ceros
        if col_motriz1:
            df['Motriz_1'] = pd.to_numeric(df[col_motriz1], errors='coerce').fillna(0).astype(int)
        else:
            df['Motriz_1'] = pd.Series(0, index=df.index, dtype=int)
        
        # Motriz 2: similar
        if col_motriz2:
            df['Motriz_2'] = pd.to_numeric(df[col_motriz2], errors='coerce').fillna(0).astype(int)
        else:
            df['Motriz_2'] = pd.Series(0, index=df.index, dtype=int)
        
        # Unidad original (si existe)
        if col_unidad:
            df['Unidad_Original'] = df[col_unidad]
        else:
            df['Unidad_Original'] = ''
        
        # Calcular Unidad según Motriz 2: 0 -> S, >0 -> M
        df['Unidad'] = df['Motriz_2'].apply(lambda x: 'M' if x > 0 else 'S')
        
        # Identificar primera estación (Puerto) y última (Limache) por sus nombres
        puerto_col = None
        limache_col = None
        for col in df.columns:
            if 'puerto' in col.lower() and 'hora salida' in col.lower():
                puerto_col = col
            if 'limache' in col.lower() and 'hora llegada' in col.lower():
                limache_col = col
        
        # Calcular hora real de salida (primera estación)
        if puerto_col:
            df['Hora_Salida_Real'] = df[puerto_col].apply(convertir_a_minutos)
        else:
            df['Hora_Salida_Real'] = None
            st.warning(f"No se encontró columna de salida (Puerto) en el archivo {file.name}. Los tiempos de salida se mostrarán como vacíos.")
        
        # Calcular hora real de llegada (última estación)
        if limache_col:
            df['Hora_Llegada_Real'] = df[limache_col].apply(convertir_a_minutos)
        else:
            df['Hora_Llegada_Real'] = None
            st.warning(f"No se encontró columna de llegada (Limache) en el archivo {file.name}. El TDV se mostrará como vacío.")
        
        # Hora programada
        df['Min_Prog'] = df['Hora_Prog'].apply(convertir_a_minutos)
        
        # Retraso (PS)
        df['Retraso'] = df['Hora_Salida_Real'] - df['Min_Prog']
        df['Puntual'] = (abs(df['Retraso']) <= 5).astype(int)
        
        # TDV
        if puerto_col and limache_col:
            tdv = df['Hora_Llegada_Real'] - df['Hora_Salida_Real']
            tdv = tdv.apply(lambda x: x if x > 0 else x + 1440)
            df['TDV_Min'] = tdv
        else:
            df['TDV_Min'] = 0
        
        # Determinar tipo de recorrido (origen-destino)
        def cod_estacion(nombre):
            if 'puerto' in nombre.lower():
                return 'PU'
            elif 'limache' in nombre.lower():
                return 'LI'
            elif 'vina' in nombre.lower() or 'viña' in nombre.lower():
                return 'VM'
            elif 'belloto' in nombre.lower():
                return 'EB'
            else:
                return nombre[:2]
        if puerto_col and limache_col:
            origen = cod_estacion(puerto_col)
            destino = cod_estacion(limache_col)
            df['Tipo_Rec'] = f"{origen}-{destino}"
        else:
            df['Tipo_Rec'] = 'OTRO'
        
        # Distancia base y Tren-Km
        df['Dist_Base'] = df['Tipo_Rec'].map(DISTANCIAS).fillna(0)
        df['Peso'] = df['Unidad'].apply(lambda x: 2 if x == 'M' else 1)
        df['Tren-Km'] = df['Dist_Base'] * df['Peso']
        
        # Fecha
        fecha_info = leer_fecha_archivo(file)
        if fecha_info:
            df['Fecha_Op'] = f"{fecha_info[0]:02d}/{fecha_info[1]:02d}/{fecha_info[2]}"
        else:
            df['Fecha_Op'] = ''
        
        # Renombrar Motriz_1 y Motriz_2 a los nombres finales deseados
        df.rename(columns={'Motriz_1': 'Motriz 1', 'Motriz_2': 'Motriz 2'}, inplace=True)
        
        # Eliminar columnas duplicadas (si las hay)
        df = df.loc[:, ~df.columns.duplicated()]
        
        # Asegurar columnas necesarias
        for col in ['Servicio', 'Motriz 1', 'Motriz 2', 'Unidad', 'Tipo_Rec', 'Tren-Km', 'Retraso', 'Puntual', 'Hora_Prog', 'Fecha_Op']:
            if col not in df.columns:
                df[col] = 0 if col in ['Servicio', 'Motriz 1', 'Motriz 2', 'Tren-Km'] else ''
        
        return df, df['Tren-Km'].sum(), df[df['TDV_Min'] > 0]['TDV_Min'].mean(), (df['Puntual'].sum() / len(df) * 100) if len(df) > 0 else 0
    except Exception as e:
        st.error(f"Error procesando THDR: {e}")
        return pd.DataFrame(), 0, 0, 0

# --- 4. INICIALIZACIÓN DE DATAFRAMES VACÍOS ---
df_ops = pd.DataFrame()
df_tr = pd.DataFrame()
df_tr_acum = pd.DataFrame()
df_seat = pd.DataFrame()
df_energy_master = pd.DataFrame()
df_p_d = pd.DataFrame()
df_f_d = pd.DataFrame()
df_thdr_v1 = pd.DataFrame()
df_thdr_v2 = pd.DataFrame()
all_comp_full = []
all_prmte_15 = []
all_fact_h = []

# --- 5. INTERFAZ DE USUARIO (SIDEBAR) ---
with st.sidebar:
    st.header("📅 Filtro Global")
    today = date.today()
    start_of_month = today.replace(day=1) if today.day > 1 else (today.replace(month=today.month-1, day=1) if today.month>1 else today.replace(year=today.year-1, month=12, day=1))
    date_range = st.date_input("Selecciona el período", value=(start_of_month, today))
    start_date, end_date = (date_range[0], date_range[1]) if isinstance(date_range, tuple) and len(date_range)==2 else (date_range, date_range)
    st.divider()
    st.header("📂 Carga de Archivos")
    f_v1 = st.file_uploader("1. THDR Vía 1", type=["xls", "xlsx"], accept_multiple_files=True)
    f_v2 = st.file_uploader("2. THDR Vía 2", type=["xls", "xlsx"], accept_multiple_files=True)
    f_umr = st.file_uploader("3. UMR / Odómetros", type=["xlsx"], accept_multiple_files=True)
    f_seat_files = st.file_uploader("4. Energía SEAT", type=["xlsx"], accept_multiple_files=True)
    f_bill_files = st.file_uploader("5. Facturación y PRMTE", type=["xlsx"], accept_multiple_files=True)

# --- 6. LECTURA Y PROCESAMIENTO DE DATOS (UMR, SEAT, PRMTE, THDR) ---
if f_v1 or f_v2 or f_umr or f_seat_files or f_bill_files:
    all_ops, all_tr, all_tr_acum, all_seat, all_prmte_15, all_fact_h, all_comp_full = [], [], [], [], [], [], []
    thdr_v1_list = []
    thdr_v2_list = []
    todos = (f_v1 or []) + (f_v2 or []) + (f_umr or []) + (f_seat_files or []) + (f_bill_files or [])

    for f in todos:
        try:
            xl = pd.ExcelFile(f)
            for sn in xl.sheet_names:
                sn_up = sn.upper()
                # UMR / RESUMEN
                if any(k in sn_up for k in ['UMR', 'RESUMEN']):
                    df_raw = pd.read_excel(f, sheet_name=sn, header=None)
                    h_r = next((i for i in range(min(100, len(df_raw))) if any(k in str(df_raw.iloc[i]).upper() for k in ['ODO', 'FECHA'])), None)
                    if h_r is not None:
                        df_p = pd.read_excel(f, sheet_name=sn, header=h_r)
                        df_p.columns = [re.sub(r'[^A-Z]', '', str(c).upper().replace('Ó','O')) for c in df_p.columns]
                        idx_f, idx_o, idx_t = next((c for c in df_p.columns if 'FECHA' in c), None), next((c for c in df_p.columns if 'ODO' in c and 'ACUM' not in c), None), next((c for c in df_p.columns if 'TREN' in c and 'KM' in c), None)
                        if idx_f and idx_o:
                            df_p['_dt'] = pd.to_datetime(df_p[idx_f], errors='coerce')
                            mask = (df_p['_dt'].dt.date >= start_date) & (df_p['_dt'].dt.date <= end_date)
                            for _, r in df_p[mask].dropna(subset=['_dt']).iterrows():
                                all_ops.append({"Fecha": r['_dt'].normalize(), "Tipo Día": get_tipo_dia(r['_dt']), "N° Semana": r['_dt'].isocalendar()[1], "Odómetro [km]": parse_latam_number(r[idx_o]), "Tren-Km [km]": parse_latam_number(r[idx_t]), "UMR [%]": (parse_latam_number(r[idx_t])/parse_latam_number(r[idx_o])*100 if parse_latam_number(r[idx_o])>0 else 0)})
                # ODOMETRO/KILOMETRAJE
                if 'ODO' in sn_up and 'KIL' in sn_up:
                    df_tr_raw = pd.read_excel(f, sheet_name=sn, header=None)
                    headers_found = []
                    for i in range(len(df_tr_raw)-2):
                        for j in range(1, len(df_tr_raw.columns)):
                            val = pd.to_datetime(df_tr_raw.iloc[i, j], errors='coerce')
                            if pd.notna(val) and start_date <= val.date() <= end_date:
                                if i not in [h[0] for h in headers_found]: headers_found.append((i, val))
                    for idx, (row_idx, s_dt) in enumerate(headers_found):
                        is_acum = any(k in str(df_tr_raw.iloc[row_idx:row_idx+3, 0:5]).upper() for k in ['ACUM', 'LECTURA', 'TOTAL'])
                        c_map = {j: pd.to_datetime(df_tr_raw.iloc[row_idx, j], errors='coerce') for j in range(1, len(df_tr_raw.columns)) if pd.notna(pd.to_datetime(df_tr_raw.iloc[row_idx, j], errors='coerce'))}
                        for k in range(row_idx+3, min(row_idx+40, len(df_tr_raw))):
                            n_tr = str(df_tr_raw.iloc[k, 0]).strip().upper()
                            if re.match(r'^(M|XM)', n_tr):
                                for c_idx, c_fch in c_map.items():
                                    val_km = parse_latam_number(df_tr_raw.iloc[k, c_idx])
                                    d_pt = {"Tren": n_tr, "Fecha": c_fch.normalize(), "Día": c_fch.day, "Valor": val_km}
                                    if is_acum or idx > 0: all_tr_acum.append(d_pt)
                                    else: all_tr.append(d_pt)
                # SEAT
                if 'SEAT' in sn_up and 'SER' in sn_up:
                    df_s = pd.read_excel(f, sheet_name=sn, header=None)
                    for i in range(len(df_s)):
                        fs = pd.to_datetime(df_s.iloc[i, 1], errors='coerce')
                        if pd.notna(fs) and start_date <= fs.date() <= end_date:
                            tot, tra, k12 = parse_latam_number(df_s.iloc[i, 3]), parse_latam_number(df_s.iloc[i, 5]), parse_latam_number(df_s.iloc[i, 7])
                            all_seat.append({"Fecha": fs.normalize(), "Total [kWh]": tot, "Tracción [kWh]": tra, "12 KV [kWh]": k12, "% Tracción": (tra/tot*100 if tot>0 else 0), "% 12 KV": (k12/tot*100 if tot>0 else 0)})
                # PRMTE / MEDIDAS
                if any(k in sn_up for k in ['PRMTE', 'MEDIDAS']):
                    df_prm = pd.read_excel(f, sheet_name=sn, header=None)
                    h_idx = next((i for i in range(len(df_prm)) if 'AÑO' in str(df_prm.iloc[i]).upper()), None)
                    if h_idx is not None:
                        df_pd = pd.read_excel(f, sheet_name=sn, header=h_idx)
                        df_pd['Timestamp'] = pd.to_datetime(df_pd[['AÑO', 'MES', 'DIA', 'HORA']].astype(int).rename(columns={'AÑO':'year','MES':'month','DIA':'day','HORA':'hour'})) + pd.to_timedelta(df_pd['INICIO INTERVALO'].astype(int), unit='m')
                        cols_e = [c for c in df_pd.columns if 'Retiro_Energia_Activa (kWhD)' in str(c)]
                        for _, r in df_pd.iterrows():
                            ts, val_p = r['Timestamp'], sum([parse_latam_number(r[col]) for col in cols_e])
                            all_comp_full.append({"Fecha": ts.normalize(), "Hora": ts.hour, "Consumo Horario [kWh]": val_p, "Fuente": "PRMTE"})
                            if start_date <= ts.date() <= end_date: all_prmte_15.append({"Fecha y Hora": ts.strftime('%d/%m/%Y %H:%M'), "Fecha": ts.normalize(), "Energía PRMTE [kWh]": val_p})
                # FACTURA / CONSUMO
                if any(k in sn_up for k in ['FACTURA', 'CONSUMO']):
                    df_f = pd.read_excel(f, sheet_name=sn); df_f.columns = ['FechaHora', 'Valor']
                    df_f['Timestamp'] = pd.to_datetime(df_f['FechaHora'], errors='coerce')
                    for _, r in df_f.dropna(subset=['Timestamp']).iterrows():
                        ts, val_f = r['Timestamp'], abs(parse_latam_number(r['Valor']))
                        all_comp_full.append({"Fecha": ts.normalize(), "Hora": ts.hour, "Consumo Horario [kWh]": val_f, "Fuente": "Factura"})
                        if start_date <= ts.date() <= end_date: all_fact_h.append({"Fecha y Hora": ts.strftime('%d/%m/%Y %H:%M'), "Fecha": ts.normalize(), "Consumo Horario [kWh]": val_f})
        except: continue

    # Procesar THDR Vía 1
    if f_v1:
        for file in f_v1:
            df, _, _, _ = procesar_thdr_avanzado(file)
            if not df.empty:
                df['Vía'] = 'Vía 1'
                thdr_v1_list.append(df)
        if thdr_v1_list:
            df_thdr_v1 = pd.concat(thdr_v1_list, ignore_index=True)
    # Procesar THDR Vía 2
    if f_v2:
        for file in f_v2:
            df, _, _, _ = procesar_thdr_avanzado(file)
            if not df.empty:
                df['Vía'] = 'Vía 2'
                thdr_v2_list.append(df)
        if thdr_v2_list:
            df_thdr_v2 = pd.concat(thdr_v2_list, ignore_index=True)

    # Jerarquía y pre‑filtrado para otras pestañas
    if any([all_ops, all_tr, all_tr_acum, all_seat, all_prmte_15, all_fact_h]):
        if all_ops: df_ops = pd.DataFrame(all_ops).drop_duplicates(subset=['Fecha']).sort_values("Fecha")
        if all_tr: df_tr = pd.DataFrame(all_tr).sort_values(["Fecha", "Tren"])
        if all_tr_acum: df_tr_acum = pd.DataFrame(all_tr_acum).sort_values(["Fecha", "Tren"])
        if all_seat: df_seat = pd.DataFrame(all_seat).drop_duplicates(subset=['Fecha']).sort_values("Fecha")
        
        if not df_seat.empty:
            df_energy_master = df_seat[["Fecha", "Total [kWh]", "Tracción [kWh]", "12 KV [kWh]"]].copy().rename(columns={"Total [kWh]":"E_Total", "Tracción [kWh]":"E_Tr", "12 KV [kWh]":"E_12"})
            df_energy_master["Fuente"] = "SEAT"

        if all_prmte_15:
            df_p_d = pd.DataFrame(all_prmte_15).groupby("Fecha")["Energía PRMTE [kWh]"].sum().reset_index()
            if not df_seat.empty:
                df_p_d = pd.merge(df_p_d, df_seat[["Fecha", "% Tracción", "% 12 KV"]], on="Fecha", how="left").fillna(0)
                df_p_d["E_Tr"], df_p_d["E_12"] = df_p_d["Energía PRMTE [kWh]"]*(df_p_d["% Tracción"]/100), df_p_d["Energía PRMTE [kWh]"]*(df_p_d["% 12 KV"]/100)
                df_p_p = df_p_d.rename(columns={"Energía PRMTE [kWh]":"E_Total"})[["Fecha","E_Total","E_Tr","E_12"]]; df_p_p["Fuente"] = "PRMTE"
                df_energy_master = pd.concat([df_energy_master, df_p_p]).drop_duplicates(subset=["Fecha"], keep="last")

        if all_fact_h:
            df_f_d = pd.DataFrame(all_fact_h).groupby("Fecha")["Consumo Horario [kWh]"].sum().reset_index()
            if not df_seat.empty:
                df_f_d = pd.merge(df_f_d, df_seat[["Fecha", "% Tracción", "% 12 KV"]], on="Fecha", how="left").fillna(0)
                df_f_d["E_Tr"], df_f_d["E_12"] = df_f_d["Consumo Horario [kWh]"]*(df_f_d["% Tracción"]/100), df_f_d["Consumo Horario [kWh]"]*(df_f_d["% 12 KV"]/100)
                df_f_f = df_f_d.rename(columns={"Consumo Horario [kWh]":"E_Total"})[["Fecha","E_Total","E_Tr","E_12"]]; df_f_f["Fuente"] = "Factura"
                df_energy_master = pd.concat([df_energy_master, df_f_f]).drop_duplicates(subset=["Fecha"], keep="last")

        if not df_ops.empty and not df_energy_master.empty:
            df_ops = pd.merge(df_ops, df_energy_master, on="Fecha", how="left")
            df_ops['IDE (kWh/km)'] = df_ops.apply(lambda row: row['E_Tr'] / row['Odómetro [km]'] if row['Odómetro [km]'] > 0 else 0, axis=1)

# --- 7. DASHBOARD ---
tabs = st.tabs(["📊 Resumen", "📑 Operaciones", "📑 Trenes", "⚡ Energía", "⚖️ Comparación Energía hr", "📈 Regresión Nocturna", "🚨 Datos Atípicos", "📋 THDR"])

# ================== PESTAÑA RESUMEN ==================
with tabs[0]:
    st.header("📊 Resumen General de Gestión Energética")
    if df_ops.empty:
        st.info("No hay datos cargados o el período seleccionado no contiene información. Sube archivos en la barra lateral.")
    else:
        # Calcular métricas principales del período
        total_energia = df_ops['E_Tr'].sum() if 'E_Tr' in df_ops else 0
        total_odometro = df_ops['Odómetro [km]'].sum() if 'Odómetro [km]' in df_ops else 0
        total_tren_km = df_ops['Tren-Km [km]'].sum() if 'Tren-Km [km]' in df_ops else 0
        ide_promedio = (total_energia / total_odometro) if total_odometro > 0 else 0
        umr_promedio = df_ops['UMR [%]'].mean() if 'UMR [%]' in df_ops else 0
        
        col1, col2, col3, col4 = st.columns(4)
        col1.metric("⚡ Energía Tracción", f"{total_energia:,.0f} kWh")
        col2.metric("🚆 Odómetro Total", f"{total_odometro:,.0f} km")
        col3.metric("📊 IDE Promedio", f"{ide_promedio:.2f} kWh/km")
        col4.metric("📈 UMR Promedio", f"{umr_promedio:.1f} %")
        
        # Gráfico de evolución diaria
        st.subheader("Evolución Diaria de IDE y UMR")
        df_plot = df_ops[['Fecha', 'IDE (kWh/km)', 'UMR [%]']].copy().dropna()
        if not df_plot.empty:
            fig = make_subplots(specs=[[{"secondary_y": True}]])
            fig.add_trace(go.Scatter(x=df_plot['Fecha'], y=df_plot['IDE (kWh/km)'], name="IDE (kWh/km)", line=dict(color='blue')), secondary_y=False)
            fig.add_trace(go.Scatter(x=df_plot['Fecha'], y=df_plot['UMR [%]'], name="UMR (%)", line=dict(color='red')), secondary_y=True)
            fig.update_layout(title_text="IDE y UMR por día", xaxis_title="Fecha", height=500)
            fig.update_yaxes(title_text="IDE (kWh/km)", secondary_y=False)
            fig.update_yaxes(title_text="UMR (%)", secondary_y=True)
            st.plotly_chart(fig, use_container_width=True)
        else:
            st.info("No hay datos suficientes para graficar IDE y UMR.")
        
        # Distribución de energía por fuente (si hay múltiples fuentes)
        if 'Fuente' in df_energy_master.columns:
            st.subheader("Energía Total por Fuente")
            energia_fuente = df_energy_master.groupby('Fuente')['E_Total'].sum().reset_index()
            fig_pie = go.Figure(data=[go.Pie(labels=energia_fuente['Fuente'], values=energia_fuente['E_Total'], hole=0.4)])
            fig_pie.update_layout(title_text="Participación de fuentes de energía")
            st.plotly_chart(fig_pie, use_container_width=True)

# ================== PESTAÑA OPERACIONES ==================
with tabs[1]:
    st.header("📑 Operaciones Diarias")
    if df_ops.empty:
        st.info("No hay datos de operaciones para el período seleccionado.")
    else:
        st.dataframe(df_ops.style.format({
            'Odómetro [km]': '{:,.0f}',
            'Tren-Km [km]': '{:,.0f}',
            'UMR [%]': '{:.1f}',
            'E_Tr': '{:,.0f}',
            'IDE (kWh/km)': '{:.2f}'
        }), use_container_width=True)
        csv_ops = df_ops.to_csv(index=False).encode('utf-8')
        st.download_button("Descargar CSV", csv_ops, "operaciones.csv", "text/csv")

# ================== PESTAÑA TRENES ==================
with tabs[2]:
    st.header("📑 Kilometraje por Tren")
    if not df_tr.empty:
        st.subheader("Kilometraje diario por tren (no acumulado)")
        st.dataframe(df_tr, use_container_width=True)
    if not df_tr_acum.empty:
        st.subheader("Lecturas acumuladas de odómetros por tren")
        st.dataframe(df_tr_acum, use_container_width=True)
    if df_tr.empty and df_tr_acum.empty:
        st.info("No hay datos de kilometraje de trenes.")

# ================== PESTAÑA ENERGÍA ==================
with tabs[3]:
    st.header("⚡ Energía SEAT")
    if df_seat.empty:
        st.info("No hay datos de energía SEAT para el período seleccionado.")
    else:
        st.dataframe(df_seat.style.format({
            'Total [kWh]': '{:,.0f}',
            'Tracción [kWh]': '{:,.0f}',
            '12 KV [kWh]': '{:,.0f}',
            '% Tracción': '{:.1f}',
            '% 12 KV': '{:.1f}'
        }), use_container_width=True)
        fig_energy = go.Figure()
        fig_energy.add_trace(go.Bar(x=df_seat['Fecha'], y=df_seat['Tracción [kWh]'], name='Tracción'))
        fig_energy.add_trace(go.Bar(x=df_seat['Fecha'], y=df_seat['12 KV [kWh]'], name='12 KV'))
        fig_energy.update_layout(title="Consumo energético diario", barmode='stack', xaxis_title="Fecha", yaxis_title="kWh")
        st.plotly_chart(fig_energy, use_container_width=True)

# ================== PESTAÑA COMPARACIÓN ENERGÍA HR ==================
with tabs[4]:
    st.header("⚖️ Comparación de Consumo Horario (PRMTE vs Factura)")
    if not all_comp_full:
        st.info("No hay datos de consumo horario (PRMTE o Factura).")
    else:
        df_comp = pd.DataFrame(all_comp_full)
        df_comp = df_comp[(df_comp['Fecha'].dt.date >= start_date) & (df_comp['Fecha'].dt.date <= end_date)]
        if df_comp.empty:
            st.info("No hay datos en el rango seleccionado.")
        else:
            pivot = df_comp.pivot_table(index=['Fecha', 'Hora'], columns='Fuente', values='Consumo Horario [kWh]', aggfunc='sum').reset_index()
            st.dataframe(pivot, use_container_width=True)
            fig_comp = go.Figure()
            for fuente in ['PRMTE', 'Factura']:
                if fuente in pivot.columns:
                    df_filt = pivot[['Fecha', 'Hora', fuente]].dropna()
                    fig_comp.add_trace(go.Scatter(x=df_filt['Hora'], y=df_filt[fuente], mode='lines+markers', name=fuente))
            fig_comp.update_layout(title="Perfil horario de consumo", xaxis_title="Hora del día", yaxis_title="kWh")
            st.plotly_chart(fig_comp, use_container_width=True)

# ================== PESTAÑA REGRESIÓN NOCTURNA ==================
with tabs[5]:
    st.header("📈 Regresión Lineal Nocturna (IDE vs Odómetro)")
    st.markdown("Relación entre energía de tracción y kilometraje recorrido para servicios nocturnos (puede ajustarse según criterio).")
    if 'E_Tr' in df_ops.columns and 'Odómetro [km]' in df_ops.columns:
        df_night = df_ops.dropna(subset=['E_Tr', 'Odómetro [km]'])
        if len(df_night) > 1:
            slope, intercept, r_value, p_value, std_err = stats.linregress(df_night['Odómetro [km]'], df_night['E_Tr'])
            fig_reg = go.Figure()
            fig_reg.add_trace(go.Scatter(x=df_night['Odómetro [km]'], y=df_night['E_Tr'], mode='markers', name='Datos'))
            x_line = np.linspace(df_night['Odómetro [km]'].min(), df_night['Odómetro [km]'].max(), 100)
            y_line = intercept + slope * x_line
            fig_reg.add_trace(go.Scatter(x=x_line, y=y_line, mode='lines', name=f'Regresión: E = {slope:.2f} * km + {intercept:.0f}'))
            fig_reg.update_layout(title=f"Regresión Lineal (R² = {r_value**2:.3f})", xaxis_title="Odómetro [km]", yaxis_title="Energía Tracción [kWh]")
            st.plotly_chart(fig_reg, use_container_width=True)
            st.metric("Pendiente (IDE promedio)", f"{slope:.2f} kWh/km")
        else:
            st.warning("Se necesitan al menos dos puntos para realizar regresión.")
    else:
        st.info("Datos insuficientes para regresión.")

# ================== PESTAÑA DATOS ATÍPICOS ==================
with tabs[6]:
    st.header("🚨 Detección de Datos Atípicos (IDE)")
    if 'IDE (kWh/km)' in df_ops.columns:
        df_ide = df_ops[['Fecha', 'IDE (kWh/km)']].dropna()
        if not df_ide.empty:
            q1 = df_ide['IDE (kWh/km)'].quantile(0.25)
            q3 = df_ide['IDE (kWh/km)'].quantile(0.75)
            iqr = q3 - q1
            lower = q1 - 1.5 * iqr
            upper = q3 + 1.5 * iqr
            outliers = df_ide[(df_ide['IDE (kWh/km)'] < lower) | (df_ide['IDE (kWh/km)'] > upper)]
            st.write(f"**Rango esperado (IQR):** {lower:.2f} - {upper:.2f} kWh/km")
            st.write(f"**Días atípicos encontrados:** {len(outliers)}")
            if not outliers.empty:
                st.dataframe(outliers)
            else:
                st.success("No se detectaron valores atípicos.")
        else:
            st.info("No hay datos de IDE.")
    else:
        st.info("No existe columna IDE para analizar.")

# ================== PESTAÑA THDR ==================
with tabs[7]:
    st.header("📋 Datos THDR - Vía 1 y Vía 2")
    
    with st.expander("🔍 Ver estructura de los DataFrames THDR (depuración)"):
        if not df_thdr_v1.empty:
            st.write("**Columnas en THDR Vía 1:**", list(df_thdr_v1.columns))
            st.dataframe(df_thdr_v1.head())
        else:
            st.info("No hay datos para Vía 1.")
        if not df_thdr_v2.empty:
            st.write("**Columnas en THDR Vía 2:**", list(df_thdr_v2.columns))
            st.dataframe(df_thdr_v2.head())
        else:
            st.info("No hay datos para Vía 2.")
    
    def mostrar_tabla_thdr(df, titulo, color_emoji):
        st.subheader(f"{color_emoji} {titulo}")
        if df.empty:
            st.info(f"No hay datos de THDR para {titulo}. Sube archivos en la sección correspondiente.")
            return
        
        df_calc = df.copy()
        # Formatear columnas calculadas
        if 'Hora_Salida_Real' in df_calc.columns:
            df_calc['Hora Real Salida'] = df_calc['Hora_Salida_Real'].apply(lambda x: format_hms(x) if pd.notna(x) else "")
        else:
            df_calc['Hora Real Salida'] = ""
        if 'Min_Prog' in df_calc.columns:
            df_calc['Hora Programada'] = df_calc['Min_Prog'].apply(lambda x: format_hms(x) if pd.notna(x) else "")
        else:
            df_calc['Hora Programada'] = ""
        if 'Retraso' in df_calc.columns:
            df_calc['Puntualidad Salida'] = df_calc['Retraso'].apply(lambda x: format_hms(x, con_signo=True) if pd.notna(x) else "")
        else:
            df_calc['Puntualidad Salida'] = ""
        if 'TDV_Min' in df_calc.columns:
            df_calc['TDV'] = df_calc['TDV_Min'].apply(lambda x: format_hms(x) if pd.notna(x) else "")
        else:
            df_calc['TDV'] = ""
        if 'Hora_Llegada_Real' in df_calc.columns:
            df_calc['Hora Llegada Terminal'] = df_calc['Hora_Llegada_Real'].apply(lambda x: format_hms(x) if pd.notna(x) else "")
        else:
            df_calc['Hora Llegada Terminal'] = ""
        
        columnas_mostrar = {
            'Fecha': 'Fecha_Op',
            'Servicio': 'Servicio',
            'Hora Programada': 'Hora Programada',
            'Hora Real Salida': 'Hora Real Salida',
            'Puntualidad Salida': 'Puntualidad Salida',
            'Hora Llegada Terminal': 'Hora Llegada Terminal',
            'TDV': 'TDV',
            'Motriz 1': 'Motriz 1',
            'Motriz 2': 'Motriz 2',
            'Unidad': 'Unidad',
            'Recorrido': 'Tipo_Rec',
            'Tren-Km': 'Tren-Km'
        }
        df_display = pd.DataFrame()
        for nombre, col_orig in columnas_mostrar.items():
            if col_orig in df_calc.columns:
                df_display[nombre] = df_calc[col_orig]
            elif nombre == 'Fecha' and 'Fecha_Op' not in df_calc.columns and 'Fecha' in df_calc.columns:
                df_display[nombre] = df_calc['Fecha']
        if df_display.empty:
            st.warning("No se encontraron columnas esperadas. Mostrando datos crudos:")
            st.dataframe(df)
            return
        if 'Tren-Km' in df_display.columns:
            df_display['Tren-Km'] = df_display['Tren-Km'].apply(lambda x: f"{x:,.1f}" if isinstance(x, (int, float)) else x)
        st.dataframe(df_display, use_container_width=True)
    
    mostrar_tabla_thdr(df_thdr_v1, "Vía 1", "🟢")
    mostrar_tabla_thdr(df_thdr_v2, "Vía 2", "🔵")

# --- 8. DESCARGA DE REPORTE EXCEL COMPLETO ---
st.sidebar.download_button("📥 Reporte Excel Completo", to_excel_consolidado(df_ops, df_tr, df_tr_acum, df_seat, df_p_d, pd.DataFrame(all_prmte_15), pd.DataFrame(all_fact_h), df_f_d), "Reporte_EFE_SGE.xlsx")
